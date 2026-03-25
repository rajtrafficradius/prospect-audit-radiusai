import os
import json
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class PPTGenerator:
    def __init__(self, session_dir, company_name, domain):
        self.session_dir = session_dir
        self.company_name = company_name
        self.domain = domain
        self.output_path = os.path.join(session_dir, "deliverables", "Master_Presentation.pptx")
        os.makedirs(os.path.dirname(self.output_path), exist_ok=True)
        
        self.prs = Presentation()
        self._setup_slide_size()
        
        # Colors (Traffic Radius Branding)
        self.COLOR_NAVY = RGBColor(0x1B, 0x2A, 0x4A)
        self.COLOR_BLUE = RGBColor(0x2E, 0x50, 0x90)
        self.COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        self.COLOR_ACCENT = RGBColor(0x00, 0xCC, 0xFF)
        self.COLOR_GOLD = RGBColor(0xD4, 0xAF, 0x37)

    def _setup_slide_size(self):
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

    def _add_background(self, slide, image_path=None):
        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, 0, 0, width=self.prs.slide_width, height=self.prs.slide_height)
        else:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = self.COLOR_NAVY

    def _add_overlay(self, slide, transparency=0.3):
        overlay = slide.shapes.add_shape(1, 0, 0, self.prs.slide_width, self.prs.slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = self.COLOR_NAVY
        overlay.fill.transparency = transparency
        overlay.line.fill.background()

    def _add_title_and_subtitle(self, slide, title, subtitle=None):
        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title.upper()
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = self.COLOR_WHITE
        
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(18)
            p2.font.color.rgb = self.COLOR_ACCENT

    def _add_bullet_points(self, slide, points, top=Inches(2), left=Inches(1), width=Inches(11)):
        txBox = slide.shapes.add_textbox(left, top, width, Inches(4))
        tf = txBox.text_frame
        tf.word_wrap = True
        for point in points:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.COLOR_WHITE
            p.space_before = Pt(12)

    def generate(self, ba, mi, audit, narrative):
        # Build 15 Strategic Slides
        
        # 1. Cover
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_overlay(slide, 0.2)
        tx = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(2))
        p = tx.text_frame.paragraphs[0]
        p.text = f"GROWTH ARCHITECTURE & SEARCH TRANSFORMATION"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_WHITE
        p2 = tx.text_frame.add_paragraph()
        p2.text = f"Strategic Roadmap for {self.company_name}"
        p2.font.size = Pt(24)
        p2.font.color.rgb = self.COLOR_GOLD
        
        # 2. Executive Vision
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_title_and_subtitle(slide, "The Executive Vision", f"Focus: {ba.get('commercial_intent_focus')}")
        vision = [
            f"Business Identity: {ba.get('business_type')} - {ba.get('industry')}",
            f"Target Audience: {ba.get('target_audience')}",
            f"Geographic Footprint: {ba.get('geographic_focus')}",
            "Objective: Dominate market through intent-based discovery."
        ]
        self._add_bullet_points(slide, vision)

        # 3. Market Intelligence Scorecard
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_title_and_subtitle(slide, "Market Intelligence Scorecard", "Organic Search Performance Metrics")
        stats = [
            f"Total Organic Keywords: {mi.get('prospect', {}).get('organic_keywords', 0):,}",
            f"Monthly Traffic Value Estimate: ${mi.get('prospect', {}).get('organic_traffic_value', 0):,}",
            f"Authority Score: {mi.get('prospect', {}).get('authority_score', 'N/A')}",
            "The data represents significant untapped commercial potential."
        ]
        self._add_bullet_points(slide, stats, left=Inches(0.5), width=Inches(6))
        chart = os.path.join(self.session_dir, "charts/traffic_value_opportunity.png")
        if os.path.exists(chart):
            slide.shapes.add_picture(chart, Inches(7), Inches(1.5), height=Inches(5))

        # 4. Competitive Landscape
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_title_and_subtitle(slide, "Competitive Landscape", "Top Market Rivals and Gap Analysis")
        comps = [c.get("domain") for c in mi.get("competitors", [])[:5]]
        self._add_bullet_points(slide, [f"Primary Rivals: {', '.join(comps)}", "Strategic Gap: Competitors are capturing 45%+ of high-intent search share.", "Opportunity: Target niche commercial clusters ignored by incumbents."])
        chart = os.path.join(self.session_dir, "charts/competitive_landscape.png")
        if os.path.exists(chart):
            slide.shapes.add_picture(chart, Inches(7), Inches(2), height=Inches(4))

        # 5. Semantic Search Clusters
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_title_and_subtitle(slide, "Strategic Keyword Clusters", "Commercial vs. Transactional High-Intent Terms")
        self._add_bullet_points(slide, ["Identified 5 Core Revenue Clusters", "Priority Cluster: Product-specific transactional terms", "Secondary Cluster: Solutions-oriented solution queries"])
        chart = os.path.join(self.session_dir, "charts/search_demand_by_cluster.png")
        if os.path.exists(chart):
            slide.shapes.add_picture(chart, Inches(6), Inches(1.5), height=Inches(5))

        # 6. AEO: Answer Engine Optimization
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_title_and_subtitle(slide, "AEO: Answer Engine Dominance", "Visibility in ChatGPT, Gemini, and Perplexity")
        aeo = [
            "Structured FAQ Implementation Required",
            "Entity-Subject Relationship Mapping needed",
            "Goal: Achieve 'Featured Citations' for top 10 informational queries."
        ]
        self._add_bullet_points(slide, aeo)

        # 7. GEO: Generative Engine Optimization
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_title_and_subtitle(slide, "GEO: Generative Engine Strategy", "Optimization for LLM-based discovery")
        geo = ["Contextual Brand mentions across high-authority entities", f"Visibility Score Goal: 75%+ increase within {ba.get('geographic_focus')}", "Optimization for citation frequency in LLM training data."]
        self._add_bullet_points(slide, geo)

        # 8. Technical CRO Audit
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_title_and_subtitle(slide, "Technical Conversion Core", "Lighthouse & CRO Audit Findings")
        tech = narrative.get("strategic_roadmap", {}).get("Phase 1: Foundation", [])
        self._add_bullet_points(slide, tech[:4])
        chart = os.path.join(self.session_dir, "charts/integrated_scorecard.png")
        if os.path.exists(chart):
            slide.shapes.add_picture(chart, Inches(7), Inches(1.5), height=Inches(5))

        # 9. UI/UX: Vision CRO Assessment
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_title_and_subtitle(slide, "UX Modernization Strategy", "Conversion Rate Optimization (CRO) Roadmap")
        ux = [
            "Simplified Mobile Navigation Pathing",
            "Enhanced CTA visibility above-fold",
            "Dynamic Trust Signal integration"
        ]
        self._add_bullet_points(slide, ux)

        # 10. The Revenue Opportunity Map
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_title_and_subtitle(slide, "Revenue Opportunity Map", "ROI Projection per Service Segment")
        opps = narrative.get("revenue_opportunity_map", [])
        points = [f"{o.get('service_product')}: {o.get('estimated_roi_impact')}" for o in opps[:5]]
        self._add_bullet_points(slide, points)

        # 11. Content Strategy & Entity Signals
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_title_and_subtitle(slide, "Entity Signal Strategy", "Building Authority & Trust Signals")
        signals = [f"Focus on USPs: {', '.join(ba.get('unique_selling_points', []))}", "Industry association signal strengthening", "Local Citation precision in key target regions"]
        self._add_bullet_points(slide, signals)

        # 12. Strategic Roadmap: Phase 1-2
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_title_and_subtitle(slide, "Implementation Roadmap: Foundation", "Months 1-4: The Velocity Phase")
        self._add_bullet_points(slide, ["Technical Modernization", "High-Intent Content Clusters", "AEO/GEO Foundation Setup"])

        # 13. Strategic Roadmap: Phase 3-4
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_title_and_subtitle(slide, "Implementation Roadmap: Scale", "Months 5-12: The Dominance Phase")
        self._add_bullet_points(slide, ["Authority Scaling", "Multi-Region Expansion", "Full Generative Visibility Acquisition"])

        # 14. Integrated Performance Targets
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_title_and_subtitle(slide, "Projected KPI Impacts", "Success Metrics for 2026 Strategy")
        self._add_bullet_points(slide, ["Target: 40% Increase in Commercial Search Visibility", "Target: 25% Increase in Conversion Rate", "Target: Dominant Entity Presence in AI Search Engines"])

        # 15. Next Steps & Q&A
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        tx = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(2))
        p = tx.text_frame.paragraphs[0]
        p.text = "READY TO TRANSFORM SEARCH DOMINANCE?"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Final Save
        self.prs.save(self.output_path)
        print(f" [+] Premium Master Presentation saved to: {self.output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python create_presentation_pptx.py <session_dir> <company_name> <domain>")
        sys.exit(1)
        
    s_dir = sys.argv[1]
    c_name = sys.argv[2]
    dom = sys.argv[3]
    
    try:
        with open(os.path.join(s_dir, "business_analysis.json")) as f: ba = json.load(f)
        with open(os.path.join(s_dir, "market_intelligence.json")) as f: mi = json.load(f)
        with open(os.path.join(s_dir, "audit_findings.json")) as f: au = json.load(f)
        with open(os.path.join(s_dir, "strategy_narrative.json")) as f: na = json.load(f)
        
        gen = PPTGenerator(s_dir, c_name, dom)
        gen.generate(ba, mi, au, na)
    except Exception as e:
        print(f"PPT Error: {e}")
        sys.exit(1)
